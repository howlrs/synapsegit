#!/usr/bin/env python3
"""Generate the Japanese SynapseGit intended-user scenario deck."""

from __future__ import annotations

import argparse
import re
from pathlib import Path
from typing import Iterable, Sequence

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


DEFAULT_OUTPUT = Path(__file__).with_name("synapsegit_user_scenarios_ja.pptx")

MAIN = "https://github.com/howlrs/synapsegit/blob/main"
LINKS = {
    "readme": f"{MAIN}/README.md",
    "usage": f"{MAIN}/docs/usage_guide.md",
    "core": f"{MAIN}/docs/core_concept.md",
    "stage0": f"{MAIN}/docs/stage0_execution_plan.md",
    "runtime": f"{MAIN}/docs/runtime_architecture.md",
    "protocol": f"{MAIN}/spec/core/v0.1/README.md",
    "presentation": f"{MAIN}/docs/presentations/README.md",
}

FONT = "Noto Sans JP"
SLIDE_W = 13.333
SLIDE_H = 7.5


def color(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


NAVY = color("0B1324")
INK = color("13202B")
MUTED = color("5E6873")
PAPER = color("F7F4ED")
WHITE = color("FFFFFF")
HAIRLINE = color("7B858A")
DARK_LINE = color("71839E")

EVIDENCE = color("0B6E69")
EVIDENCE_SOFT = color("D9EEEB")
DECISION = color("8C4A12")
DECISION_SOFT = color("F3E3C7")
PROPOSAL = color("74428A")
PROPOSAL_SOFT = color("E8E1F3")
ANALYSIS = color("5E6873")
ANALYSIS_SOFT = color("E7EAED")
GAP = color("A13A31")
GAP_SOFT = color("F5DEDB")
ACTIVITY = color("8C450E")
ACTIVITY_SOFT = color("F5E4D4")
PLAN = color("4E51A8")
PLAN_SOFT = color("E2E4F7")
SUCCESS = color("29704A")
SUCCESS_SOFT = color("DCEEE4")
CYAN = color("3A8EA8")


def i(value: float):
    return Inches(value)


def set_background(slide, fill: RGBColor) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = fill


def set_alt_text(shape, title: str, description: str) -> None:
    """Set the non-visual properties used by PowerPoint accessibility tools."""
    try:
        nodes = shape._element.xpath(".//p:cNvPr")
        if nodes:
            clear_decorative(shape)
            nodes[0].set("name", title)
            nodes[0].set("title", title)
            nodes[0].set("descr", description)
    except Exception:
        # Alt text is helpful metadata, but never a generation blocker.
        pass


DECORATIVE_NS = "http://schemas.microsoft.com/office/drawing/2017/decorative"
DECORATIVE_URI = "{C183D7F6-B498-43B3-948B-1728B52AA6E4}"


def set_decorative(shape) -> None:
    try:
        c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
        ext_list = c_nv_pr.find(qn("a:extLst"))
        if ext_list is None:
            ext_list = OxmlElement("a:extLst")
            c_nv_pr.append(ext_list)
        if any(node.tag == f"{{{DECORATIVE_NS}}}decorative" for node in ext_list.iter()):
            return
        extension = OxmlElement("a:ext")
        extension.set("uri", DECORATIVE_URI)
        decorative = etree.Element(
            f"{{{DECORATIVE_NS}}}decorative", nsmap={"adec": DECORATIVE_NS}
        )
        decorative.set("val", "1")
        extension.append(decorative)
        ext_list.append(extension)
    except Exception:
        pass


def clear_decorative(shape) -> None:
    try:
        c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
        ext_list = c_nv_pr.find(qn("a:extLst"))
        if ext_list is None:
            return
        for extension in list(ext_list):
            if any(
                node.tag == f"{{{DECORATIVE_NS}}}decorative"
                for node in extension.iter()
            ):
                ext_list.remove(extension)
        if len(ext_list) == 0:
            c_nv_pr.remove(ext_list)
    except Exception:
        pass


def style_run(run, size: float, fill: RGBColor, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = fill
    properties = run._r.get_or_add_rPr()
    properties.set("lang", "ja-JP")
    east_asian = properties.find(qn("a:ea"))
    if east_asian is None:
        east_asian = OxmlElement("a:ea")
        properties.append(east_asian)
    east_asian.set("typeface", FONT)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 18,
    fill: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
    name: str | None = None,
):
    box = slide.shapes.add_textbox(i(x), i(y), i(w), i(h))
    box.name = name or f"Text: {text[:24]}"
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = i(margin)
    frame.margin_right = i(margin)
    frame.margin_top = i(margin)
    frame.margin_bottom = i(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    for index, line in enumerate(text.split("\n")):
        if index:
            paragraph.add_line_break()
        run = paragraph.add_run()
        run.text = line
        style_run(run, size, fill, bold)
    return box


def add_rich_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    parts: Sequence[tuple[str, RGBColor, bool]],
    *,
    size: float = 18,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
):
    box = slide.shapes.add_textbox(i(x), i(y), i(w), i(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = i(0.03)
    frame.margin_right = i(0.03)
    frame.margin_top = i(0.02)
    frame.margin_bottom = i(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, fill, bold in parts:
        run = paragraph.add_run()
        run.text = text
        style_run(run, size, fill, bold)
    return box


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    items: Iterable[str],
    *,
    size: float = 15,
    fill: RGBColor = INK,
    bullet_fill: RGBColor | None = None,
    spacing: float = 5,
):
    box = slide.shapes.add_textbox(i(x), i(y), i(w), i(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = i(0.02)
    frame.margin_right = i(0.02)
    frame.margin_top = i(0.02)
    frame.margin_bottom = i(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(spacing)
        bullet = paragraph.add_run()
        bullet.text = "● "
        style_run(bullet, max(8, size - 3), bullet_fill or fill, True)
        run = paragraph.add_run()
        run.text = item
        style_run(run, size, fill, False)
    return box


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    line_width: float = 1,
    name: str | None = None,
):
    shape = slide.shapes.add_shape(shape_type, i(x), i(y), i(w), i(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    if name:
        shape.name = name
    set_decorative(shape)
    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    fill: RGBColor = HAIRLINE,
    width: float = 1.5,
    dashed: bool = False,
):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, i(x1), i(y1), i(x2), i(y2)
    )
    connector.line.color.rgb = fill
    connector.line.width = Pt(width)
    if dashed:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    set_decorative(connector)
    return connector


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    title: str,
    body: str,
    accent: RGBColor,
    soft: RGBColor,
    tag: str | None = None,
    title_size: float = 18,
    body_size: float = 14,
):
    card = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill=WHITE,
        line=HAIRLINE,
        line_width=0.8,
        name=f"Card: {title}",
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x + 0.18,
        y + 0.18,
        0.14,
        h - 0.36,
        fill=accent,
        name=f"Accent: {title}",
    )
    title_size = max(title_size, 15.5)
    body_size = max(body_size, 14.0)
    if tag:
        add_pill(slide, x + 0.48, y + 0.18, min(w - 0.68, 1.42), 0.34, tag, accent, soft, 10)
        title_y = y + 0.62
    else:
        title_y = y + 0.14
    add_text(slide, x + 0.48, title_y, w - 0.68, 0.38, title, size=title_size, fill=accent, bold=True)
    body_y = title_y + 0.44
    add_text(
        slide,
        x + 0.48,
        body_y,
        w - 0.68,
        max(0.34, y + h - body_y - 0.12),
        body,
        size=body_size,
        fill=INK,
    )
    return card


def add_pill(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    accent: RGBColor,
    soft: RGBColor,
    size: float = 11,
    *,
    link: str | None = None,
):
    pill = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill=soft,
        line=accent,
        line_width=0.7,
        name=f"Pill: {text}",
    )
    frame = pill.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    style_run(run, size, accent, True)
    if link:
        run.hyperlink.address = link
        pill.click_action.hyperlink.address = link
    set_alt_text(pill, text, text)
    return pill


def set_slide_title(slide, text: str, x: float, y: float, w: float, h: float, *, size: float, fill: RGBColor):
    title = slide.shapes.title
    if title is None:
        return add_text(slide, x, y, w, h, text, size=size, fill=fill, bold=True)
    title.left = i(x)
    title.top = i(y)
    title.width = i(w)
    title.height = i(h)
    title.name = f"Slide title: {text[:32]}"
    frame = title.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = i(0.02)
    frame.margin_right = i(0.02)
    frame.margin_top = i(0.02)
    frame.margin_bottom = i(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    for index, line in enumerate(text.split("\n")):
        if index:
            paragraph.add_line_break()
        run = paragraph.add_run()
        run.text = line
        style_run(run, size, fill, True)
    return title


def add_title(slide, kicker: str, title: str, subtitle: str | None = None, *, dark=False):
    main = WHITE if dark else INK
    muted = color("BDC6D3") if dark else MUTED
    add_text(slide, 0.64, 0.35, 2.9, 0.32, kicker.upper(), size=10.5, fill=CYAN if dark else EVIDENCE, bold=True)
    set_slide_title(slide, title, 0.64, 0.68, 12.0, 0.62, size=28, fill=main)
    if subtitle:
        add_text(slide, 0.66, 1.28, 11.8, 0.4, subtitle, size=15, fill=muted)


def add_footer(slide, number: int, source: str, *, dark=False):
    line = color("334057") if dark else HAIRLINE
    text = color("AAB5C5") if dark else MUTED
    add_line(slide, 0.64, 7.05, 12.68, 7.05, fill=line, width=0.7)
    box = add_text(
        slide,
        0.64,
        7.10,
        3.0,
        0.2,
        "Concept visual · branch: main",
        size=8.5,
        fill=text,
    )
    run = box.text_frame.paragraphs[0].runs[0]
    run.hyperlink.address = LINKS["usage"]
    add_text(slide, 3.75, 7.10, 7.9, 0.2, source, size=8.5, fill=text, align=PP_ALIGN.CENTER)
    add_text(slide, 12.18, 7.10, 0.5, 0.2, f"{number:02d}", size=8.5, fill=text, bold=True, align=PP_ALIGN.RIGHT)


def add_node(slide, x, y, w, h, label, sublabel, accent, soft, *, dashed=False):
    shape = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill=soft,
        line=accent,
        line_width=1.2,
        name=f"Node: {label}",
    )
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    add_text(slide, x + 0.12, y + 0.14, w - 0.24, 0.36, label, size=15, fill=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.10, y + 0.55, w - 0.20, 0.34, sublabel, size=12.5, fill=INK, align=PP_ALIGN.CENTER)
    return shape


def slide_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.16, 7.5, fill=EVIDENCE)
    add_pill(slide, 0.72, 0.52, 2.28, 0.38, "CORE v0.1 · STAGE 0", CYAN, color("142D45"), 11.5)
    set_slide_title(slide, "SynapseGit\n利用イメージ", 0.72, 1.30, 6.4, 1.15, size=37, fill=WHITE)
    add_text(
        slide,
        0.75,
        2.72,
        5.85,
        1.12,
        "完成物だけでなく、\n何を変え、なぜ採用したかを次へ渡す。",
        size=22,
        fill=color("DCE6F2"),
        bold=True,
    )
    add_text(slide, 0.75, 4.15, 5.7, 0.85, "画家・壁画家／建築家／施工・修復／\nデザイナー＋Creative AI", size=15.5, fill=color("AAB8C8"))
    add_pill(slide, 0.75, 5.52, 3.52, 0.44, "利用構想 · UIスクリーンショットではありません", WHITE, color("233047"), 10.5)

    # Branching lineage graphic.
    add_line(slide, 7.55, 1.35, 7.55, 5.95, fill=DARK_LINE, width=3)
    add_line(slide, 7.55, 2.52, 9.15, 1.80, fill=PROPOSAL, width=2)
    add_line(slide, 7.55, 3.48, 9.40, 3.48, fill=EVIDENCE, width=2)
    add_line(slide, 7.55, 4.45, 9.15, 5.18, fill=ACTIVITY, width=2)
    add_node(slide, 6.70, 0.86, 1.70, 0.92, "意図", "Plan", PLAN, PLAN_SOFT)
    add_node(slide, 6.70, 2.12, 1.70, 0.92, "実行", "Activity", ACTIVITY, ACTIVITY_SOFT)
    add_node(slide, 6.70, 3.10, 1.70, 0.92, "観測", "Observation", EVIDENCE, EVIDENCE_SOFT)
    add_node(slide, 6.70, 4.08, 1.70, 0.92, "意味付け", "Claim", PROPOSAL, PROPOSAL_SOFT)
    add_node(slide, 6.55, 5.30, 2.00, 1.00, "人が判断", "Decision Commit", DECISION, DECISION_SOFT)
    add_card(slide, 9.20, 1.25, 3.10, 1.35, title="今日", body="比較して、次の一手を決める", accent=EVIDENCE, soft=EVIDENCE_SOFT, title_size=17, body_size=13)
    add_card(slide, 9.45, 2.88, 3.10, 1.35, title="次の案件", body="手順・却下案・判断基準を再利用", accent=PLAN, soft=PLAN_SOFT, title_size=17, body_size=13)
    add_card(slide, 9.20, 4.54, 3.10, 1.35, title="将来", body="改修・修復・再展示時に根拠へ戻る", accent=DECISION, soft=DECISION_SOFT, title_size=17, body_size=13)
    add_text(slide, 8.55, 6.31, 3.95, 0.46, "作者証明・完全な現実記録・\n自動利益分配ではありません", size=11.5, fill=color("C3CDDA"), align=PP_ALIGN.RIGHT)
    add_text(slide, 0.75, 6.31, 5.95, 0.46, "Scope: SynapseGit Coreのみ\nChrono-Engine／歴史的人物再現は対象外", size=11.5, fill=color("C3CDDA"))
    add_footer(slide, 1, "Core concept §§1, 20 · Stage 0 draft", dark=True)


def slide_audiences(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, PAPER)
    add_title(slide, "WHO", "誰の、どんな断絶を解くのか", "異なる現場から、同じ「判断の系譜」へ接続する")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.66, 1.58, 12.0, 0.52, fill=NAVY, name="Shared problem")
    add_text(slide, 0.88, 1.68, 11.56, 0.28, "写真・図面・チャット・AI案が散在　→　判断理由が消える　→　次の担当が最初から探し直す", size=14.5, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
    cards = [
        ("画家・壁画家", "色・構図・偶然・保留が\n完成画像へ埋もれる", "同条件比較／制作process pack", EVIDENCE, EVIDENCE_SOFT),
        ("建築家", "Brief・案・現場判断が\n別資料へ散る", "三者比較／設計変更理由", PLAN, PLAN_SOFT),
        ("施工・修復", "計画外作業・欠測・処置が\n後任へ伝わらない", "Hold Point／引渡し記録", ACTIVITY, ACTIVITY_SOFT),
        ("デザイナー＋AI", "ツールやAIを跨ぐと\n採否の文脈が消える", "提案系列／可搬な文脈", PROPOSAL, PROPOSAL_SOFT),
    ]
    for index, (title, problem, benefit, accent, soft) in enumerate(cards):
        x = 0.66 + index * 3.03
        add_card(slide, x, 2.35, 2.79, 3.25, title=title, body=f"困りごと\n{problem}\n\n返るもの\n{benefit}", accent=accent, soft=soft, title_size=17, body_size=13.3)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.66, 5.90, 12.0, 0.72, fill=WHITE, line=HAIRLINE, line_width=0.8)
    add_rich_text(slide, 0.92, 6.02, 11.48, 0.42, [
        ("二次利用者　", MUTED, True),
        ("制作チーム・後任・施主・所有者・コレクター・美術館　", INK, False),
        ("← 報告・引き継ぎ・保存資料として受け取る", SUCCESS, True),
    ], size=13.5, align=PP_ALIGN.CENTER)
    add_text(slide, 1.30, 6.62, 10.7, 0.32, "制作ソフト、BIM/CAD、ペイントツールを置き換えるのではなく、判断を横断してつなぐ。", size=13, fill=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 2, "Core concept §§14, 20.2", dark=False)


def slide_core_loop(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, PAPER)
    add_title(slide, "HOW", "節目で残し、すぐ返す", "セッション終了・承認・引き渡し・Hold Pointだけを意味ある履歴にする")
    nodes = [
        ("意図", "Plan", PLAN, PLAN_SOFT, False),
        ("実行", "Activity", ACTIVITY, ACTIVITY_SOFT, False),
        ("観測", "Observation", EVIDENCE, EVIDENCE_SOFT, False),
        ("比較候補", "Analysis", ANALYSIS, ANALYSIS_SOFT, True),
        ("意味付け", "Claim", PROPOSAL, PROPOSAL_SOFT, False),
        ("人の判断", "Decision", DECISION, DECISION_SOFT, False),
    ]
    xs = [0.66 + index * 2.03 for index in range(6)]
    for index in range(5):
        add_shape(slide, MSO_SHAPE.CHEVRON, xs[index] + 1.66, 2.54, 0.34, 0.44, fill=HAIRLINE)
    for x, (label, sub, accent, soft, dashed) in zip(xs, nodes):
        add_node(slide, x, 2.20, 1.65, 1.12, label, sub, accent, soft, dashed=dashed)
    add_text(slide, 5.82, 3.43, 1.75, 0.32, "人が採否を判断", size=13, fill=DECISION, bold=True, align=PP_ALIGN.CENTER)

    returns = [
        ("比較", "前後・計画との違い", EVIDENCE, EVIDENCE_SOFT),
        ("記録復元", "その節目の記録を取り出す", PLAN, PLAN_SOFT),
        ("報告", "進捗・処置の下書き", ACTIVITY, ACTIVITY_SOFT),
        ("引き継ぎ", "理由・制約・未解決", SUCCESS, SUCCESS_SOFT),
    ]
    for index, (title, body, accent, soft) in enumerate(returns):
        x = 0.90 + index * 3.07
        add_card(slide, x, 4.14, 2.76, 1.42, title=title, body=body, accent=accent, soft=soft, title_size=16, body_size=12.5)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.90, 5.87, 11.56, 0.82, fill=WHITE, line=HAIRLINE, line_width=0.8)
    add_rich_text(slide, 1.12, 5.99, 11.12, 0.56, [
        ("見方の原則　", INK, True),
        ("写真＝観測　", EVIDENCE, True),
        ("差分＝再計算可能な解析　", ANALYSIS, True),
        ("確定事実ではない　", GAP, True),
        ("常時監視・隠れた録画は行わない", INK, False),
    ], size=14, align=PP_ALIGN.CENTER)
    add_footer(slide, 3, "Core concept §§3, 14.4, 20.4", dark=False)


def draw_canvas(slide, x, y, w, h, variant: str):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill=color("EFE9DE"), line=color("B6AA98"), line_width=1.2)
    add_shape(slide, MSO_SHAPE.ARC, x + 0.30, y + 0.42, 1.52, 1.20, fill=PLAN_SOFT, line=PLAN, line_width=1.2)
    add_shape(slide, MSO_SHAPE.OVAL, x + 1.45, y + 0.40, 0.74, 0.74, fill=DECISION_SOFT, line=DECISION, line_width=1.0)
    add_shape(slide, MSO_SHAPE.WAVE, x + 0.40, y + 1.25, 1.76, 0.56, fill=EVIDENCE_SOFT, line=EVIDENCE, line_width=1.0)
    if variant == "after":
        add_shape(slide, MSO_SHAPE.FREEFORM if hasattr(MSO_SHAPE, "FREEFORM") else MSO_SHAPE.CLOUD, x + 1.10, y + 0.92, 0.86, 0.74, fill=GAP_SOFT, line=GAP, line_width=2)
        add_pill(slide, x + 1.14, y + 1.50, 1.28, 0.34, "差分候補", GAP, GAP_SOFT, 11)


def slide_painter(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, PAPER)
    add_title(slide, "SCENARIO 01", "画家・壁画家 — 塗り重ねの「なぜ」を残す", "同じ視点の前後を比べ、差分候補へ人が意味を付ける")
    add_text(slide, 0.72, 1.55, 2.7, 0.30, "制作前 Capture", size=14, fill=EVIDENCE, bold=True, align=PP_ALIGN.CENTER)
    draw_canvas(slide, 0.78, 1.94, 2.58, 2.18, "before")
    add_text(slide, 4.05, 1.55, 2.4, 0.30, "制作 Session", size=14, fill=ACTIVITY, bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 4.12, 1.94, 2.28, 2.18, fill=ACTIVITY_SOFT, line=ACTIVITY, line_width=1.2)
    add_text(slide, 4.43, 2.22, 1.65, 0.52, "描く\n塗り重ねる", size=19, fill=ACTIVITY, bold=True, align=PP_ALIGN.CENTER)
    add_pill(slide, 4.34, 3.18, 1.84, 0.42, "任意の一言・音声", ACTIVITY, WHITE, 12)
    add_shape(slide, MSO_SHAPE.CHEVRON, 3.46, 2.76, 0.42, 0.52, fill=HAIRLINE)
    add_shape(slide, MSO_SHAPE.CHEVRON, 6.58, 2.76, 0.42, 0.52, fill=HAIRLINE)
    add_text(slide, 7.05, 1.55, 2.7, 0.30, "制作後 Capture", size=14, fill=EVIDENCE, bold=True, align=PP_ALIGN.CENTER)
    draw_canvas(slide, 7.10, 1.94, 2.58, 2.18, "after")
    add_card(slide, 10.05, 1.72, 2.58, 2.60, title="すぐ返るもの", body="前後比較\n差分候補\n制作process pack\n未採用案の再利用", accent=SUCCESS, soft=SUCCESS_SOFT, title_size=17, body_size=13.5)

    add_text(slide, 0.82, 4.46, 2.2, 0.34, "色・構図案を分ける", size=14, fill=PROPOSAL, bold=True)
    add_line(slide, 1.10, 5.24, 3.55, 5.24, fill=PROPOSAL, width=2)
    add_line(slide, 2.05, 5.24, 2.90, 4.83, fill=PROPOSAL, width=1.5)
    add_line(slide, 2.05, 5.24, 2.90, 5.65, fill=PROPOSAL, width=1.5)
    add_pill(slide, 3.03, 4.63, 1.60, 0.42, "採用した色案", SUCCESS, SUCCESS_SOFT, 10.5)
    add_pill(slide, 3.03, 5.48, 1.85, 0.42, "保留・未採用案", PROPOSAL, PROPOSAL_SOFT, 10.5)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.15, 4.56, 7.46, 1.42, fill=WHITE, line=HAIRLINE, line_width=0.8)
    add_rich_text(slide, 5.42, 4.78, 6.92, 0.95, [
        ("人が残す意味　", DECISION, True),
        ("「この変化を採用」「偶然を残す」「まだ不明」\n", INK, False),
        ("しないこと　", GAP, True),
        ("差分量を創造性・努力の点数にしない／画像から作者を特定しない", INK, False),
    ], size=14, valign=MSO_ANCHOR.TOP)
    add_footer(slide, 4, "Core concept §§14.2, 20.2, 20.7", dark=False)


def draw_plan_panel(slide, x, y, w, h, mode: str):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill=WHITE, line=HAIRLINE, line_width=1)
    for offset in (0.45, 0.90, 1.35, 1.80, 2.25):
        add_line(slide, x + 0.20, y + offset, x + w - 0.20, y + offset, fill=color("EDF0F1"), width=0.6)
    accent = PLAN if mode == "plan" else EVIDENCE
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.44, y + 0.48, 0.16, 1.72, fill=accent)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.44, y + 0.48, 1.92, 0.16, fill=accent)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 2.18, y + 0.48, 0.16, 1.18, fill=accent)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 1.15, y + 1.55, 1.19, 0.16, fill=accent)
    if mode == "current":
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 1.52, y + 1.72, 0.54, 0.45, fill=GAP_SOFT, line=GAP, line_width=1.6)
        add_pill(slide, x + 1.56, y + 2.16, 1.28, 0.34, "要確認", GAP, GAP_SOFT, 11)


def slide_architect(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, PAPER)
    add_title(slide, "SCENARIO 02", "建築家 — Plan・直前現況・現在現況を並べる", "「計画との適合」と「時間変化」を混ぜずに判断する")
    panels = [
        (0.68, "PLAN", "実現したかった状態", "plan", PLAN),
        (4.50, "PREVIOUS", "直前に観測した状態", "previous", EVIDENCE),
        (8.32, "CURRENT", "現在観測した状態", "current", EVIDENCE),
    ]
    for x, title, sub, mode, accent in panels:
        add_text(slide, x, 1.58, 3.34, 0.28, title, size=14, fill=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x, 1.87, 3.34, 0.30, sub, size=12.5, fill=MUTED, align=PP_ALIGN.CENTER)
        draw_plan_panel(slide, x, 2.22, 3.34, 2.80, mode)
    add_shape(slide, MSO_SHAPE.CHEVRON, 4.05, 3.25, 0.32, 0.50, fill=HAIRLINE)
    add_shape(slide, MSO_SHAPE.CHEVRON, 7.87, 3.25, 0.32, 0.50, fill=HAIRLINE)
    add_text(slide, 3.42, 4.88, 1.60, 0.30, "計画との適合", size=12.5, fill=PLAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 7.24, 4.88, 1.60, 0.30, "時間変化", size=12.5, fill=EVIDENCE, bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 11.74, 2.22, 0.94, 2.80, fill=DECISION_SOFT, line=DECISION, line_width=1.4)
    add_text(slide, 11.82, 2.50, 0.78, 2.14, "人が判断\n\n採用\n是正\n保留", size=14, fill=DECISION, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    chips = [
        ("設計変更理由", PLAN, PLAN_SOFT),
        ("施主・現場との認識合わせ", EVIDENCE, EVIDENCE_SOFT),
        ("次案件へ判断基準を再利用", SUCCESS, SUCCESS_SOFT),
    ]
    for index, (label, accent, soft) in enumerate(chips):
        add_pill(slide, 0.88 + index * 4.02, 5.55, 3.62, 0.48, label, accent, soft, 13.5)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.88, 6.23, 11.70, 0.50, fill=GAP_SOFT, line=GAP, line_width=0.8)
    add_text(slide, 1.05, 6.31, 11.36, 0.30, "BIM/CADの代替ではない。写真中心の記録は自動的にAs-built認定せず、確認範囲付きのAs-recordedとして扱う。", size=12, fill=GAP, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 5, "Core concept §§14.1, 14.3, 20.2", dark=False)


def slide_construction(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, PAPER)
    add_title(slide, "SCENARIO 03", "施工・修復 — Hold Pointを未来の担当へ渡す", "処置前後だけでなく、材料・可逆性・観測できなかった範囲も残す")
    xs = [1.10, 3.85, 6.55, 9.25, 11.80]
    add_line(slide, xs[0], 3.55, xs[-1], 3.55, fill=ACTIVITY, width=4)
    stages = [
        ("手順", "Procedure", PLAN, PLAN_SOFT),
        ("処置前", "Hold Point", DECISION, DECISION_SOFT),
        ("作業", "Activity", ACTIVITY, ACTIVITY_SOFT),
        ("処置後", "Capture / Review", EVIDENCE, EVIDENCE_SOFT),
        ("引渡し", "Pack", SUCCESS, SUCCESS_SOFT),
    ]
    for x, (label, sub, accent, soft) in zip(xs, stages):
        add_shape(slide, MSO_SHAPE.OVAL, x - 0.22, 3.33, 0.44, 0.44, fill=accent, line=WHITE, line_width=1.2)
        add_text(slide, x - 0.70, 3.86, 1.40, 0.32, label, size=14, fill=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x - 0.76, 4.18, 1.52, 0.28, sub, size=11.5, fill=MUTED, align=PP_ALIGN.CENTER)

    add_card(slide, 0.68, 1.62, 2.46, 1.33, title="予定", body="工程・完了条件\n観測要件", accent=PLAN, soft=PLAN_SOFT, title_size=16, body_size=12.5)
    add_card(slide, 3.18, 1.62, 2.46, 1.33, title="処置前 Capture", body="対象・位置・状態\n任意の音声", accent=EVIDENCE, soft=EVIDENCE_SOFT, title_size=16, body_size=12.5)
    add_card(slide, 5.74, 4.70, 2.46, 1.33, title="材料・処置", body="入力・出力\n可逆性・計画との差", accent=ACTIVITY, soft=ACTIVITY_SOFT, title_size=16, body_size=12.5)
    add_card(slide, 8.30, 1.62, 2.46, 1.33, title="Coverage", body="見えた範囲\n遮蔽・品質警告", accent=EVIDENCE, soft=EVIDENCE_SOFT, title_size=16, body_size=12.5)
    add_card(slide, 10.23, 4.70, 2.46, 1.33, title="成果", body="進捗・処置報告\nAs-recorded引渡し", accent=SUCCESS, soft=SUCCESS_SOFT, title_size=16, body_size=12.5)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.78, 4.78, 4.28, 1.18, fill=GAP_SOFT, line=GAP, line_width=1.0)
    add_text(slide, 1.02, 4.94, 3.80, 0.32, "EvidenceGap — 撮れなかったことも記録", size=14, fill=GAP, bold=True)
    add_text(slide, 1.02, 5.34, 3.80, 0.44, "欠測・遮蔽・緊急対応を\n「変化なし」へ置き換えない", size=13, fill=INK)
    add_text(slide, 1.05, 6.34, 11.25, 0.36, "契約適合や施工品質を自動証明しない。後任が判断できるEvidenceと、誰が何を確認したかを渡す。", size=13, fill=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 6, "Core concept §§6, 14.1, 14.4, 20.2", dark=False)


def slide_ai(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, NAVY)
    add_title(slide, "SCENARIO 04", "デザイナー＋AI — AIは枝を作る、人が採用版を決める", "文脈と制約は渡す。採用・公開・引き渡しの権限は渡さない。", dark=True)
    add_card(slide, 0.72, 1.82, 2.55, 2.18, title="Context Pack", body="Brief／Evidence\nPolicy／委任範囲\nbase commit", accent=PROPOSAL, soft=PROPOSAL_SOFT, title_size=17, body_size=13.5)
    add_shape(slide, MSO_SHAPE.CHEVRON, 3.44, 2.61, 0.42, 0.52, fill=DARK_LINE)

    # Proposal-only branches.
    add_line(slide, 4.05, 2.89, 6.20, 1.95, fill=PROPOSAL, width=2)
    add_line(slide, 4.05, 2.89, 6.20, 2.89, fill=PROPOSAL, width=2)
    add_line(slide, 4.05, 2.89, 6.20, 3.83, fill=PROPOSAL, width=2)
    for index, y in enumerate((1.58, 2.52, 3.46), start=1):
        add_pill(slide, 5.80, y, 2.05, 0.60, f"AI Proposal {index}", PROPOSAL, PROPOSAL_SOFT, 13)
    add_text(slide, 4.20, 4.30, 3.45, 0.34, "proposal/{agent}/{run} のみ", size=13, fill=color("E1D6EE"), bold=True, align=PP_ALIGN.CENTER)

    # Human gate and official history.
    gate = add_shape(slide, MSO_SHAPE.HEXAGON, 8.26, 2.06, 1.70, 1.70, fill=DECISION_SOFT, line=DECISION, line_width=3, name="Human review gate")
    add_text(slide, 8.50, 2.47, 1.22, 0.68, "HUMAN\nGATE", size=15, fill=DECISION, bold=True, align=PP_ALIGN.CENTER)
    set_alt_text(gate, "Human Gate", "人がAI提案を採用・修正・却下する")
    add_shape(slide, MSO_SHAPE.CHEVRON, 7.88, 2.62, 0.34, 0.48, fill=DARK_LINE)
    add_shape(slide, MSO_SHAPE.CHEVRON, 10.08, 2.62, 0.34, 0.48, fill=DARK_LINE)
    add_card(slide, 10.54, 1.98, 2.08, 1.86, title="公式判断", body="採用／部分採用\n修正／却下\nDecision Feedback", accent=SUCCESS, soft=SUCCESS_SOFT, title_size=16, body_size=12)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 4.92, 11.90, 1.34, fill=color("172238"), line=DARK_LINE, line_width=1.0)
    boundaries = [
        "AIは decision/*・release/* を進めない",
        "作者性・貢献率を自動判定しない",
        "明示opt-inなしに外部学習へ出さない",
    ]
    for index, text in enumerate(boundaries):
        x = 0.98 + index * 3.86
        add_pill(slide, x, 5.25, 3.48, 0.56, text, color("E1D6EE"), color("2A2340"), 11.5)
    add_text(slide, 0.92, 6.45, 11.40, 0.34, "generated_by AI ／ selected・modified・approved_by human を分離する。", size=13, fill=color("C3CDDA"), align=PP_ALIGN.CENTER)
    add_footer(slide, 7, "Core concept §21 · Stage 0 Workstream D", dark=True)


def slide_handoff(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, PAPER)
    add_title(slide, "HANDOFF", "制作チーム・後任 — 重要な変更へ2分で到達する", "完成物、チャット、フォルダを巡回せず、判断から根拠へ遡る（Pilot目標）")
    add_pill(slide, 10.73, 0.54, 1.92, 0.46, "PILOT目標 · 2分", DECISION, DECISION_SOFT, 10.5)

    cx, cy = 6.66, 3.58
    positions = [
        (0.82, 1.72, "何が変わった？", "方向付きDiff"),
        (0.82, 4.58, "なぜ変えた？", "採用理由・Claim"),
        (4.78, 5.48, "何を採らなかった？", "Proposal／却下理由"),
        (9.58, 1.72, "根拠は？", "Evidence／Observation"),
        (9.58, 4.58, "次に守ることは？", "制約・未解決事項"),
    ]
    for x, y, _, _ in positions:
        add_line(slide, cx, cy, x + 1.42, y + 0.55, fill=HAIRLINE, width=1.8)
    center = add_shape(slide, MSO_SHAPE.OVAL, 5.42, 2.45, 2.48, 2.26, fill=DECISION_SOFT, line=DECISION, line_width=2.8, name="Decision Commit")
    add_text(slide, 5.78, 2.94, 1.76, 0.78, "判断の節目\nDecision Commit", size=17, fill=DECISION, bold=True, align=PP_ALIGN.CENTER)
    set_alt_text(center, "判断の節目", "5つの問いから根拠と制約へ遡る中心点")
    accents = [EVIDENCE, DECISION, PROPOSAL, EVIDENCE, PLAN]
    softs = [EVIDENCE_SOFT, DECISION_SOFT, PROPOSAL_SOFT, EVIDENCE_SOFT, PLAN_SOFT]
    for (x, y, title, body), accent, soft in zip(positions, accents, softs):
        add_card(slide, x, y, 2.84, 1.10, title=title, body=body, accent=accent, soft=soft, title_size=14.5, body_size=11.5)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.08, 1.60, 2.50, 0.60, fill=SUCCESS_SOFT, line=SUCCESS, line_width=0.8)
    add_text(slide, 3.28, 1.69, 2.10, 0.34, "報告・引き継ぎpack", size=14, fill=SUCCESS, bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.78, 5.53, 1.62, 0.58, fill=PLAN_SOFT, line=PLAN, line_width=0.8)
    add_text(slide, 7.92, 5.61, 1.34, 0.34, "open archive", size=13, fill=PLAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.28, 6.60, 10.78, 0.34, "hashで照合できるのは記録の同一性。内容の真実性、作者性、永久保存を保証するものではない。", size=12.5, fill=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 8, "Core concept §§20.2, 20.6 · Pilot target", dark=False)


def slide_pilot(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, PAPER)
    add_title(slide, "START SMALL", "まずは「1対象 × 1節目 × 1比較」から", "絵画・壁画をprimary pilot、小規模内装・壁面を次の検証対象にする")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.18, 1.55, 10.96, 0.92, fill=NAVY, name="Pilot formula")
    add_rich_text(slide, 1.55, 1.75, 10.22, 0.50, [
        ("1対象", WHITE, True), ("　×　", color("8FBBC8"), False),
        ("1節目", WHITE, True), ("　×　", color("8FBBC8"), False),
        ("1比較", WHITE, True), ("　→　", color("8FBBC8"), False),
        ("次の判断へ戻れる状態", color("BFE2D0"), True),
    ], size=20, align=PP_ALIGN.CENTER)
    steps = [
        ("01", "対象を選ぶ", "キャンバス／壁面／一区画", PLAN, PLAN_SOFT),
        ("02", "条件を決める", "Imported／Repeatable／Calibrated", EVIDENCE, EVIDENCE_SOFT),
        ("03", "節目で撮る", "作業終了／承認／Hold Point", ACTIVITY, ACTIVITY_SOFT),
        ("04", "三者比較", "Plan／Previous／Current", ANALYSIS, ANALYSIS_SOFT),
        ("05", "人が判断", "採用／是正／保留／不明", DECISION, DECISION_SOFT),
        ("06", "渡せる形へ", "report／handoff／archive", SUCCESS, SUCCESS_SOFT),
    ]
    for index, (num, title, body, accent, soft) in enumerate(steps):
        row, col = divmod(index, 3)
        x = 0.72 + col * 4.10
        y = 2.80 + row * 1.40
        add_card(slide, x, y, 3.78, 1.12, title=f"{num}  {title}", body=body, accent=accent, soft=soft, title_size=15.5, body_size=11.5)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 5.78, 12.00, 0.76, fill=WHITE, line=HAIRLINE, line_width=0.8)
    metrics = [
        ("20秒以内", "Capture能動入力中央値"),
        ("30秒以内", "meaningful Commit"),
        ("2分以内", "重要変更の説明・引継ぎ"),
        ("100%", "空store restore成功"),
    ]
    for index, (value, label) in enumerate(metrics):
        x = 0.92 + index * 2.96
        add_text(slide, x, 5.91, 1.06, 0.28, value, size=15, fill=DECISION if index < 3 else SUCCESS, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 1.05, 5.86, 1.76, 0.40, label, size=11.5, fill=MUTED, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, 4.36, 6.58, 4.66, 0.32, "すべてPilot UX／受入目標。実績値ではありません。", size=12, fill=GAP, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 9, "Core concept §15 · Stage 0 benefit hypotheses", dark=False)


def add_link_button(slide, x, y, w, text, url, accent):
    button = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.46, fill=WHITE, line=accent, line_width=1.0, name=f"Link: {text}")
    frame = button.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    style_run(run, 12.5, accent, True)
    run.hyperlink.address = url
    button.click_action.hyperlink.address = url
    set_alt_text(button, text, f"branch main link: {url}")
    return button


def slide_current_state(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, PAPER)
    add_title(slide, "CURRENT STATE", "現在地・信頼できる境界・次の実装", "前9枚は利用構想を含みます。製品UIの完成状態を示すものではありません。")
    implemented = [
        ("実装済み · Core", "strict JSON／schema\ncanonical OID\nvalidated ingest", SUCCESS, SUCCESS_SOFT),
        ("実装済み · Local", "filesystem／Ref CAS\narchive round trip\nSQLite projection／lineage", EVIDENCE, EVIDENCE_SOFT),
        ("実装済み · App routes", "AI: preflight／Executor\nHuman: admitted handle／permit\nACL/profile FIFO fence", PROPOSAL, PROPOSAL_SOFT),
    ]
    for index, (title, body, accent, soft) in enumerate(implemented):
        x = 0.72 + index * 4.10
        add_card(slide, x, 1.62, 3.78, 1.48, title=title, body=body, accent=accent, soft=soft, title_size=15.5, body_size=12.2)

    add_card(slide, 0.72, 3.40, 5.82, 2.05, title="信頼できる見方", body="hashはbyte identityを照合\nAI／Human appはprocess-local\nHuman認証はpublish冒頭1回\nTTLが外部revocation差を限定", accent=EVIDENCE, soft=EVIDENCE_SOFT, title_size=17, body_size=13)
    add_card(slide, 6.82, 3.40, 5.90, 2.05, title="未実装・次", body="HTTP／JWT／durable ACL・permit\nOS sandbox／release／quorum\nProjection app route\nObservation／Surreal／creator UI", accent=ACTIVITY, soft=ACTIVITY_SOFT, title_size=17, body_size=13)

    add_text(slide, 0.84, 5.55, 11.64, 0.34, "しない約束：作者性・現実・契約適合の自動証明／常時監視／個人生産性score／無断AI学習／「永久保存」の過大表示", size=11.5, fill=GAP, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 0.74, 5.91, 2.10, 0.28, "branch: main resources", size=12, fill=MUTED, bold=True)
    buttons = [
        ("Usage guide", LINKS["usage"], EVIDENCE),
        ("Core concept", LINKS["core"], PLAN),
        ("Stage 0 plan", LINKS["stage0"], ACTIVITY),
        ("Protocol v0.1", LINKS["protocol"], PROPOSAL),
    ]
    for index, (label, url, accent) in enumerate(buttons):
        add_link_button(slide, 0.72 + index * 3.02, 6.20, 2.76, label, url, accent)
    add_text(slide, 0.76, 6.72, 11.92, 0.24, "リンク先はGitHub mainブランチ。main反映後に有効になります。", size=10.5, fill=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 10, "README · Security model · Stage 0 execution plan", dark=False)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = i(SLIDE_W)
    prs.slide_height = i(SLIDE_H)
    prs.core_properties.title = "SynapseGit 利用イメージ — 想定利用者別シナリオ"
    prs.core_properties.subject = "Core v0.1 / Stage 0 intended-user visualization"
    prs.core_properties.author = "SynapseGit project"
    prs.core_properties.keywords = "SynapseGit, creator, observation, decision, Creative AI"
    prs.core_properties.comments = "Generated from docs/presentations/generate_user_scenarios_pptx.py"

    slide_cover(prs)
    slide_audiences(prs)
    slide_core_loop(prs)
    slide_painter(prs)
    slide_architect(prs)
    slide_construction(prs)
    slide_ai(prs)
    slide_handoff(prs)
    slide_pilot(prs)
    slide_current_state(prs)
    return prs


def validate_presentation(path: Path) -> list[str]:
    prs = Presentation(path)
    errors: list[str] = []
    if len(prs.slides) != 10:
        errors.append(f"expected 10 slides, found {len(prs.slides)}")
    if prs.slide_width != i(SLIDE_W) or prs.slide_height != i(SLIDE_H):
        errors.append("slide size is not 13.333 × 7.5 inches")

    hyperlink_count = 0
    hyperlink_targets: set[str] = set()
    japanese = re.compile(r"[\u3040-\u30ff\u3400-\u9fff]")
    for index, slide in enumerate(prs.slides, start=1):
        if slide.shapes.title is None or not slide.shapes.title.text.strip():
            errors.append(f"slide {index} has no semantic title placeholder")
        visible_text = "\n".join(
            shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
        )
        if not visible_text.strip():
            errors.append(f"slide {index} contains no visible text")
        for relationship in slide.part.rels.values():
            if relationship.reltype.endswith("/hyperlink"):
                hyperlink_count += 1
                hyperlink_targets.add(relationship.target_ref)
                if not relationship.target_ref.startswith(f"{MAIN}/"):
                    errors.append(
                        f"slide {index} has unexpected external link {relationship.target_ref}"
                    )
        decorative_count = 0
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                errors.append(f"slide {index} has a shape outside the top/left boundary")
            if shape.left + shape.width > prs.slide_width + i(0.01):
                errors.append(f"slide {index} has a shape outside the right boundary")
            if shape.top + shape.height > prs.slide_height + i(0.01):
                errors.append(f"slide {index} has a shape outside the bottom boundary")
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                if shape.width < i(0.40) or shape.height < i(0.18):
                    errors.append(
                        f"slide {index} has an unsafe text frame {shape.name} "
                        f"({shape.width / 914400:.2f} × {shape.height / 914400:.2f} in)"
                    )
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if not run.text.strip():
                            continue
                        if shape.top < i(7.0) and (
                            run.font.size is None or run.font.size.pt < 10.5
                        ):
                            errors.append(
                                f"slide {index} has primary text below 10.5 pt in {shape.name}"
                            )
                        if japanese.search(run.text):
                            properties = run._r.get_or_add_rPr()
                            east_asian = properties.find(qn("a:ea"))
                            if properties.get("lang") != "ja-JP":
                                errors.append(
                                    f"slide {index} has Japanese text without ja-JP language"
                                )
                            if (
                                east_asian is None
                                or east_asian.get("typeface") != FONT
                            ):
                                errors.append(
                                    f"slide {index} has Japanese text without {FONT} East Asian font"
                                )
            if any(
                node.tag == f"{{{DECORATIVE_NS}}}decorative"
                for node in shape._element.iter()
            ):
                decorative_count += 1
        if decorative_count == 0:
            errors.append(f"slide {index} has no decorative accessibility markers")
    if hyperlink_count < 8:
        errors.append(f"expected at least 8 main-branch hyperlinks, found {hyperlink_count}")
    for required in (LINKS["usage"], LINKS["core"], LINKS["stage0"], LINKS["protocol"]):
        if required not in hyperlink_targets:
            errors.append(f"missing required main-branch hyperlink {required}")
    return errors


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument(
        "--check",
        action="store_true",
        help="validate the existing output instead of regenerating it",
    )
    args = parser.parse_args()
    output = args.output.resolve()

    if not args.check:
        output.parent.mkdir(parents=True, exist_ok=True)
        build_presentation().save(output)

    if not output.exists():
        raise SystemExit(f"presentation does not exist: {output}")
    errors = validate_presentation(output)
    if errors:
        for error in errors:
            print(f"error: {error}")
        return 1
    print(
        f"ok: {output} "
        "(10 slides, layout, Japanese fonts, accessibility metadata, and main links validated)"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
